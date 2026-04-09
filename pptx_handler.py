from pptx import *
import os
from pathlib import Path
import sys
import tkinter as tk
from tkinter import messagebox

def fuente_direccion(dir):
    if getattr(sys, 'frozen', False):
        base_dir = sys._MEIPASS
        
    else:
        base_dir = os.path.dirname(os.path.abspath(__file__))
    return os.path.join(base_dir, dir)
    
    
    
    # dir_base = sys._MEIPASS
    # return os.path.join(dir_base , dir)    


        
def generar_power(colonia, lugar, fecha, hora, ruta, mensaje, acompañado):
    colonia = colonia.capitalize()
    lugar = lugar.capitalize()
    fecha = fecha.capitalize()
    hora = hora.capitalize()
    mensaje = mensaje.capitalize()
    acompañado = acompañado.capitalize()
    
    ruta_formato = fuente_direccion("formato_3.pptx")
    prs = Presentation(ruta_formato)

    def get_value(filename="log.txt"):
        with open(filename, "a+") as f:
            f.seek(0)
            val = int(f.read() or 0) + 1
            f.seek(0)
            f.truncate()
            f.write(str(val))
            return val

    veces = 1


    if prs:
        print("inicio")
    else:
        print('no se encontró el formato')
        
    slide = prs.slides[0]
            
    def replace_text(shape,new_text):
        if not shape.has_text_frame:
            return
        
        
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.runs:
                paragraph.runs[0].text = new_text
                for run in paragraph.runs[1:]:
                    run.text = ""
                break


    


    for shape in slide.shapes:

        if shape.name == 'colonia':
            replace_text(shape, colonia)                                  
        elif shape.name == 'lugar':
            replace_text(shape, lugar)
        elif shape.name == 'hora':
            replace_text(shape, hora)
        elif shape.name == 'fecha':
            replace_text(shape,fecha)
        elif shape.name == 'mensaje':
            replace_text(shape, mensaje)
        elif shape.name == 'acompaña':
            replace_text(shape, acompañado)
    contador = get_value()
    # nombre_archivo = f"{contador}.pptx"
    # if nombre_archivo == "":
        # print('ERR:: NOMBRE_ARCHIVO FALLÓ')
    # archivo_salida = os.path.join(os.getcwd(), nombre_archivo)
    
    
    prs.save(ruta)
    os.startfile(ruta)
    return ruta

    