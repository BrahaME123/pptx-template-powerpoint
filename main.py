from pptx import *
import os
prs = Presentation("formato_3.pptx")
veces = 1
def get_value(filename="log.txt"):
    with open(filename, "a+") as f:
        f.seek(0)
        val = int(f.read() or 0) + 1
        f.seek(0)
        f.truncate()
        f.write(str(val))
        return val




"""
font_style = text_frame.paragraphs[0].runs[0].font
text_frame.clear()
run = text_frame.paragraphs[0].add_run()
run.font = font_style

"""

if prs:
    print("inicio")
else:
    print('no se encontró el formato')
    
slide = prs.slides[0]

"""""
             for paragraph in shape.text_frame.paragraphs:                 for run in paragraph.runs:                     run.text = "Hola"                     
             print(run.text)
"""

    
def replace_text(shape,new_text):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.runs:
            paragraph.runs[0].text = new_text
            for run in paragraph.runs[1:]:
                run.text = ""
            break


    
    
    # for paragraph in shape.text_frame.paragraphs:
    #     for run in paragraph.runs:
    #         run.text = formato[0]


        # for paragraph in shape.text_frame.paragraphs:

        #     for run in paragraph.runs:
        #         shape.text_frame.clear()
        #         run.text = formato[1]                
        #         #for run in paragraph.runs:             
        #         # run.text = formato[1]
        #         # print(run.text)
        


formato = ["Victor david", "alex tobias", "cristian suñiga", "piter es joto"]

for shape in slide.shapes:

    if shape.name == 'colonia':
        replace_text(shape, formato[0])
                                                                    
    elif shape.name == 'lugar':
        replace_text(shape, formato[1])
        
    elif shape.name == 'hora':
        replace_text(shape, formato[2])
    elif shape.name == 'fecha':
        replace_text(shape,formato[3])


    
contador = get_value()
print(f"guardado: {contador}" )
prs.save(f"{contador}.pptx")